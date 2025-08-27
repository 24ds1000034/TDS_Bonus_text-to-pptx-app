# ppt_builder.py

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.shapes import PP_PLACEHOLDER
import io
import random

# Map our layout hints to likely layout names in real templates.
LAYOUT_PREFS = [
    ("title_and_content", ["Title and Content", "Content with Caption", "Title and Content (2)", "Title and Content 2"]),
    ("title_only",        ["Title Only", "Blank Title", "Title"]),
    ("section_header",    ["Section Header", "Section Title", "Title Slide"]),
    ("two_content",       ["Two Content", "Two Content and Title", "Comparison"]),
    ("quote",             ["Quote", "Title Only"]),
    ("comparison",        ["Comparison", "Two Content"]),
    ("timeline",          ["Title and Content", "Two Content"]),
    ("process",           ["Title and Content", "Two Content"]),
    ("overview",          ["Title and Content", "Title Only"]),
    ("summary",           ["Title and Content", "Title Only"]),
]

def _find_layout_index(prs: Presentation, layout_hint: str) -> int:
    """Pick a slide layout index from the template using our hint -> name mapping."""
    prefs = next((p[1] for p in LAYOUT_PREFS if p[0] == layout_hint), None) or ["Title and Content", "Title Only"]
    # Try preferred names
    for name in prefs:
        for i, layout in enumerate(prs.slide_layouts):
            try:
                lname = getattr(layout, "name", "") or ""
                if name.lower() in lname.lower():
                    return i
            except Exception:
                pass
    # Fallback: any layout containing 'title'
    for i, layout in enumerate(prs.slide_layouts):
        try:
            lname = getattr(layout, "name", "") or ""
            if "title" in lname.lower():
                return i
        except Exception:
            pass
    # Last resort: first layout
    return 0

def _collect_template_images(prs: Presentation):
    """Collect (blob, width, height) of images present in the uploaded template/presentation."""
    images = []
    for slide in prs.slides:
        for shape in slide.shapes:
            try:
                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE and hasattr(shape, "image"):
                    images.append((shape.image.blob, shape.width, shape.height))
            except Exception:
                continue
    random.shuffle(images)
    return images

def _purge_all_existing_slides(prs: Presentation):
    """
    Remove ALL existing slides while preserving the template's masters/theme.
    Works for both PPTX and POTX (POTX typically has no slides anyway).
    """
    sldIdLst = prs.slides._sldIdLst  # oxml element
    slide_ids = list(sldIdLst)       # copy; we'll remove elements while iterating
    for sldId in slide_ids:
        rId = sldId.rId
        prs.part.drop_rel(rId)
        sldIdLst.remove(sldId)

def _first_placeholder(slide, *types):
    """Return the first placeholder whose placeholder_format.type is in types."""
    for shp in slide.placeholders:
        try:
            if shp.placeholder_format and shp.placeholder_format.type in types:
                return shp
        except Exception:
            pass
    return None

def _set_text(shape, text: str):
    """Safely set text into a placeholder or textbox if it has a text_frame."""
    try:
        if not hasattr(shape, "text_frame") or shape.text_frame is None:
            return
        tf = shape.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = text or ""
        # Gentle size reduction for very long titles
        if len(text or "") > 70:
            for r in p.runs:
                r.font.size = Pt(20)
    except Exception:
        pass

def _set_bullets(shape, bullets):
    """Write bullets into a placeholder or textbox (if it has a text_frame)."""
    try:
        if not hasattr(shape, "text_frame") or shape.text_frame is None:
            return
        tf = shape.text_frame
        tf.clear()
        bullets = bullets or []
        for i, b in enumerate(bullets[:12]):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            run = p.add_run()
            run.text = b
            p.level = 0
            if len(b) > 80:
                for r in p.runs:
                    r.font.size = Pt(16)
    except Exception:
        pass

def _ensure_notes(slide, text: str):
    """Add speaker notes if requested; ignore failures if template lacks notes master."""
    if not text:
        return
    try:
        notes = slide.notes_slide  # may create a notes slide if missing
        if notes and notes.notes_text_frame:
            notes.notes_text_frame.text = text
    except Exception:
        pass

def build_presentation(template_bytes: bytes, slides_plan):
    """
    Build a new PPTX from the uploaded template and an LLM-produced slide plan.
    slides_plan is a list of dicts: {title, bullets, layout_hint, notes?}
    """
    prs = Presentation(io.BytesIO(template_bytes))

    # 1) Collect images from the original slides, then purge slides to avoid content leakage.
    template_images = _collect_template_images(prs)
    _purge_all_existing_slides(prs)

    # 2) Build slides from plan
    for idx, slide_data in enumerate(slides_plan):
        layout_idx = _find_layout_index(prs, slide_data.get("layout_hint", "title_and_content"))
        slide = prs.slides.add_slide(prs.slide_layouts[layout_idx])

        title_text = (slide_data.get("title") or "")[:200]
        bullets = (slide_data.get("bullets") or [])[:12]

        # Prefer placeholders by role (more robust than relying on index)
        title_ph = _first_placeholder(slide,
                                      PP_PLACEHOLDER.TITLE,
                                      PP_PLACEHOLDER.CENTER_TITLE,
                                      PP_PLACEHOLDER.SUBTITLE)
        body_ph  = _first_placeholder(slide,
                                      PP_PLACEHOLDER.BODY,
                                      PP_PLACEHOLDER.CONTENT)

        if title_ph:
            _set_text(title_ph, title_text)
        else:
            try:
                tb = slide.shapes.add_textbox(Inches(1), Inches(0.7), Inches(8), Inches(1))
                _set_text(tb, title_text)
            except Exception:
                pass

        if body_ph:
            _set_bullets(body_ph, bullets)
        else:
            try:
                tb = slide.shapes.add_textbox(Inches(1), Inches(1.7), Inches(8), Inches(4.5))
                _set_bullets(tb, bullets)
            except Exception:
                pass

        # Speaker notes (optional)
        _ensure_notes(slide, slide_data.get("notes", ""))

        # Reuse a template image occasionally for visual continuity
        if template_images and ((idx % 3 == 0) or len(bullets) <= 2):
            blob, w, h = template_images[idx % len(template_images)]
            try:
                slide.shapes.add_picture(io.BytesIO(blob), Inches(0.4), Inches(5.2), height=Inches(1.2))
            except Exception:
                pass

    out = io.BytesIO()
    prs.save(out)
    return out.getvalue()
